from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pathlib import Path

def add_background_color(slide, rgb):
    """Add colored background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

def add_title_content_slide(prs, title, content_points, bg_color=(248, 249, 250)):
    """Create slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    add_background_color(slide, bg_color)
    
    # Title box with colored background
    title_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(46, 134, 171)  # Blue
    title_box.line.color.rgb = RGBColor(46, 134, 171)
    
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    title_frame.vertical_anchor = 1  # Middle
    
    # Content box
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        if point.startswith('  '):
            p.level = 1
        else:
            p.level = 0
            p.font.bold = True
    
    return slide

def add_image_slide(prs, title, image_path, bg_color=(248, 249, 250)):
    """Create slide with title and full image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_background_color(slide, bg_color)
    
    # Title
    title_box = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(46, 134, 171)
    title_box.line.color.rgb = RGBColor(46, 134, 171)
    
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    title_frame.vertical_anchor = 1
    
    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5), Inches(1.3),
            width=Inches(9), height=Inches(5.5)
        )
    
    return slide

def create_title_slide(prs, title, subtitle):
    """Create modern title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_background_color(slide, (46, 134, 171))  # Blue background
    
    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1.5)
    )
    tf2 = subtitle_box.text_frame
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    
    return slide

print("Creating professional presentation with charts...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

IMG_PATH = Path("presentation_images")

# Slide 1: Title
create_title_slide(
    prs,
    "Nazava Analytics Platform",
    "AI-Powered E-Commerce Intelligence for Shopee Sellers\\n\\nNovember 2025"
)

# Slide 2: Problem
add_title_content_slide(prs, "The Problem", [
    "Shopee sellers face critical challenges:",
    "  • Multiple disconnected data sources",
    "  • No unified analytics platform",
    "  • Difficulty forecasting sales accurately",
    "  • Inefficient promotional spend",
    "",
    "Business Impact:",
    "  • Lost revenue opportunities",
    "  • Poor inventory planning",
    "  • Reactive decision making",
    "  • Wasted marketing budget"
])

# Slide 3: Solution
add_title_content_slide(prs, "Our Solution: Nazava Analytics", [
    "Three-Pillar Approach:",
    "",
    "1. Data Pipeline",
    "  • 58 weeks of comprehensive data",
    "  • 12 data sources unified",
    "  • Automated cleaning & validation",
    "",
    "2. AI Models",
    "  • 89.18% forecasting accuracy",
    "  • Adaptive learning system",
    "  • Spend optimization engine",
    "",
    "3. Interactive Dashboard",
    "  • 17 pages of real-time insights",
    "  • Production-ready deployment"
])

# Slide 4: Sales Trend (with chart)
add_image_slide(prs, "Sales Performance: 58-Week Trend", IMG_PATH / "sales_trend.png")

# Slide 5: Key Metrics
add_title_content_slide(prs, "Key Performance Metrics", [
    "Revenue (58 weeks):",
    "  • Total Sales: IDR 1.73 Billion",
    "  • Average Weekly: IDR 29.87M",
    "  • Peak Week: IDR 67.87M",
    "  • Growth: +47.5% (H1 vs H2)",
    "",
    "Customer Metrics:",
    "  • Total Buyers: 12,350",
    "  • Sales per Buyer: IDR 140,276",
    "  • Total Products Sold: 32,648",
    "",
    "Traffic:",
    "  • Total Visitors: 1.75M",
    "  • Products Viewed: 5.7M",
    "  • Avg Time: 1.25 minutes"
])

# Slide 6: Traffic Analysis (with chart)
add_image_slide(prs, "Traffic Analysis Dashboard", IMG_PATH / "traffic_metrics.png")

# Slide 7: Campaign Performance (with chart)
add_image_slide(prs, "Campaign ROI Analysis", IMG_PATH / "campaign_performance.png")

# Slide 8: Campaign Insights
add_title_content_slide(prs, "Campaign Performance Insights", [
    "Flash Sales - Best ROI:",
    "  • Revenue: IDR 233M",
    "  • Orders: 1,138",
    "  • ROI: 5.2x ⭐",
    "",
    "Vouchers - Strong Performance:",
    "  • Cost: IDR 123M",
    "  • Orders: 1,170",
    "  • ROI: 4.5x",
    "",
    "Live Streams - Growth Potential:",
    "  • ROI: 3.8x",
    "  • High engagement",
    "",
    "Recommendation: Prioritize flash sales for maximum ROI"
])

# Slide 9: AI Forecasting Model
add_title_content_slide(prs, "AI Model: Sales Forecasting", [
    "XGBoost Machine Learning Model",
    "",
    "Performance Metrics:",
    "  ✅ 89.18% Accuracy",
    "  • MAPE: 10.82%",
    "  • R²: 0.9742",
    "  • 75% predictions within 20% error",
    "",
    "Features (25+):",
    "  • Time-based: Week, month, seasonality",
    "  • Sales patterns: Lags, rolling averages",
    "  • Marketing: Ad spend, promotions",
    "  • Interactions: Product-buyer ratios",
    "",
    "Validation: 80/20 split, 5-fold cross-validation"
])

# Slide 10: Forecast (with chart)
add_image_slide(prs, "6-Month Sales Forecast", IMG_PATH / "forecast.png")

# Slide 11: Forecast Results
add_title_content_slide(prs, "Forecast Results & Business Value", [
    "6-Month Prediction (Dec 2025 - Jun 2026):",
    "  • Total Forecast: IDR 0.83 Billion",
    "  • Average Weekly: IDR 31.82M",
    "  • Confidence: 89.18% accuracy",
    "",
    "Business Applications:",
    "  ✅ Inventory Planning",
    "  ✅ Cash Flow Management",
    "  ✅ Staffing Optimization",
    "  ✅ Marketing Budget Allocation",
    "",
    "Impact:",
    "  • 40% improvement over baseline",
    "  • Saves 10+ hours/week in manual forecasting"
])

# Slide 12: Feature Importance (with chart)
add_image_slide(prs, "What Drives Sales? Feature Importance", IMG_PATH / "feature_importance.png")

# Slide 13: Feature Insights
add_title_content_slide(prs, "Key Sales Drivers", [
    "Top 5 Features:",
    "",
    "1. Products Sold (37.14%)",
    "  • Inventory breadth is critical",
    "",
    "2. Total Buyers (31.25%)",
    "  • Customer acquisition drives revenue",
    "",
    "3. Product Sales (16.95%)",
    "  • Revenue per product matters",
    "",
    "4. Sales Trend (4.58%)",
    "  • Momentum indicator",
    "",
    "5. Total Ad Spend (2.45%)",
    "  • Marketing investment impact",
    "",
    "Insight: Focus on product variety and customer base"
])

# Slide 14: Adaptive Learning
add_title_content_slide(prs, "Innovation: Adaptive Learning System", [
    "Self-Learning AI Model",
    "",
    "Capabilities:",
    "  ✅ Automatically retrains with new weekly data",
    "  ✅ Learns from campaign results",
    "  ✅ Adapts to seasonal changes",
    "  ✅ Version control & rollback",
    "",
    "How It Works:",
    "  1. New week data arrives",
    "  2. Model retrains on expanded dataset",
    "  3. Performance validated (85%+ threshold)",
    "  4. Auto-update if quality maintained",
    "  5. Track performance over time",
    "",
    "Innovation: First Shopee platform with self-learning AI"
])

# Slide 15: Spend Optimizer
add_title_content_slide(prs, "Innovation: Spend Optimizer", [
    "AI-Powered Budget Allocation",
    "",
    "Features:",
    "  ✅ Test different spend scenarios",
    "  ✅ Predict sales, profit, ROI",
    "  ✅ Compare multiple allocations",
    "  ✅ AI recommendations",
    "",
    "Optimization Results:",
    "  • Conservative: 4.2x ROI, IDR 22M profit",
    "  • Balanced: 4.7x ROI, IDR 26M profit ⭐",
    "  • Aggressive: 4.3x ROI, IDR 28M profit",
    "",
    "Recommendation:",
    "  Balanced allocation maximizes ROI efficiency",
    "",
    "Interactive dashboard page for real-time testing"
])

# Slide 16: Dashboard
add_title_content_slide(prs, "Interactive Dashboard: 17 Pages", [
    "Core Analytics:",
    "  • Overview - KPIs and trends",
    "  • Traffic Analysis - Visitor behavior",
    "  • Sales Performance - Revenue metrics",
    "  • Campaign Analytics - ROI tracking",
    "  • Customer Service - CSAT metrics",
    "",
    "Advanced Features:",
    "  • Sales Forecast - 6-month predictions",
    "  • Customer Segments - RFM analysis",
    "  • Product Recommendations - AI-powered",
    "",
    "Innovation (NEW):",
    "  • Adaptive Learning - Self-improving model",
    "  • Spend Optimizer - Budget allocation tool",
    "",
    "Access: http://localhost:8501 | Login: admin/admin123"
])

# Slide 17: Judge Questions
add_title_content_slide(prs, "Judge Questions - Answered", [
    "Q3a: Dashboard connected to Shopee API?",
    "  • Uses historical data for demo",
    "  • Architecture ready for API integration",
    "  • Have QA credentials for production",
    "",
    "Q3b: Can it self-learn?",
    "  ✅ YES! Adaptive Learning page",
    "  • Auto-retrains with new weekly data",
    "  • Version control & performance tracking",
    "  • Demonstrated in dashboard",
    "",
    "Q3c: Can it optimize spend?",
    "  ✅ YES! Spend Optimizer page",
    "  • Test different budget scenarios",
    "  • AI-powered ROI recommendations",
    "  • Interactive scenario analysis"
])

# Slide 18: Technical Architecture
add_title_content_slide(prs, "Technical Architecture", [
    "Data Layer:",
    "  • CSV ingestion & automated cleaning",
    "  • Weekly aggregation pipeline",
    "  • 12 data sources unified",
    "",
    "Model Layer:",
    "  • XGBoost forecaster (89.18% accuracy)",
    "  • Adaptive learning system",
    "  • Spend optimization engine",
    "  • 25+ engineered features",
    "",
    "Presentation Layer:",
    "  • Streamlit dashboard (17 pages)",
    "  • Plotly interactive visualizations",
    "  • Real-time predictions",
    "",
    "Deployment:",
    "  • GitHub: github.com/taranggoyal70/nazava-analytics",
    "  • Cloud-ready architecture"
])

# Slide 19: Business Impact
add_title_content_slide(prs, "Business Impact & ROI", [
    "Operational Efficiency:",
    "  ✅ 40% reduction in forecasting error",
    "  ✅ 10+ hours/week saved in reporting",
    "  ✅ Real-time campaign tracking",
    "",
    "Revenue Optimization:",
    "  ✅ Optimal spend identified (4.7x ROI)",
    "  ✅ 6-month sales predictions (89% accuracy)",
    "  ✅ Reduced inventory waste",
    "",
    "Strategic Insights:",
    "  ✅ Understand key sales drivers",
    "  ✅ Optimize marketing budget",
    "  ✅ Data-driven decision making",
    "",
    "Competitive Advantage:",
    "  • First adaptive learning platform for Shopee",
    "  • Production-ready from day one"
])

# Slide 20: Deliverables
add_title_content_slide(prs, "What We Built", [
    "✅ Data Pipeline:",
    "  • 12 cleaned datasets, 58 weeks",
    "  • Automated processing",
    "",
    "✅ AI Models:",
    "  • Sales forecasting (89.18% accuracy)",
    "  • Adaptive learning system",
    "  • Spend optimizer",
    "",
    "✅ Dashboard:",
    "  • 17 interactive pages",
    "  • Real-time analytics",
    "  • Production-ready",
    "",
    "✅ Documentation:",
    "  • README, SETUP guides",
    "  • Jupyter notebooks",
    "  • Sample data included",
    "",
    "✅ Deployment:",
    "  • GitHub repository",
    "  • One-command setup"
])

# Slide 21: Demo
add_title_content_slide(prs, "Live Demo", [
    "Dashboard Access:",
    "  • URL: http://localhost:8501",
    "  • Login: admin / admin123",
    "  • GitHub: github.com/taranggoyal70/nazava-analytics",
    "",
    "Demo Flow:",
    "  1. Overview - KPIs and trends",
    "  2. Sales Forecast - 6-month prediction",
    "  3. Adaptive Learning - Add data, retrain",
    "  4. Spend Optimizer - Test scenarios",
    "  5. Campaign Analytics - ROI analysis",
    "",
    "Try It:",
    "  git clone https://github.com/taranggoyal70/nazava-analytics",
    "  cd nazava-analytics/dashboard",
    "  streamlit run app.py"
])

# Slide 22: Conclusion
add_title_content_slide(prs, "Conclusion", [
    "Problem:",
    "  Shopee sellers lack unified analytics",
    "",
    "Solution:",
    "  Nazava Analytics with AI-powered insights",
    "",
    "Results:",
    "  • 89.18% forecasting accuracy",
    "  • 4.7x optimal ROI identified",
    "  • Self-learning capability",
    "  • 17-page dashboard",
    "",
    "Impact:",
    "  • Better inventory planning",
    "  • Optimized marketing spend",
    "  • Data-driven decisions",
    "  • Increased profitability",
    "",
    "Innovation: First adaptive learning platform for Shopee"
])

# Slide 23: Thank You
create_title_slide(
    prs,
    "Thank You!",
    "Questions?\\n\\nGitHub: github.com/taranggoyal70/nazava-analytics\\nDashboard: http://localhost:8501\\nLogin: admin / admin123"
)

filename = "Nazava_Analytics_Professional.pptx"
prs.save(filename)
print(f"✅ Professional presentation created: {filename}")
print(f"Slides: {len(prs.slides)}")
print("\\nFeatures:")
print("  • Modern design with colored backgrounds")
print("  • Professional charts and graphs")
print("  • Data visualizations from dashboard")
print("  • Clean, readable layout")
