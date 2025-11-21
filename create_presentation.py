"""
Nazava Analytics - PowerPoint Presentation Generator
Automatically creates a professional presentation with all insights
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content_points):
    """Create bullet point slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    tf.clear()
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Create two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    return slide

def main():
    print("Creating Nazava Analytics Presentation...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "Nazava Analytics Platform",
        "AI-Powered E-Commerce Intelligence for Shopee Sellers\n\nNovember 2025"
    )
    
    # Slide 2: Problem Statement
    create_content_slide(prs, "Problem Statement", [
        "Challenge:",
        "  • Shopee sellers struggle with data-driven decision making",
        "  • Multiple disconnected data sources",
        "  • No unified analytics platform",
        "  • Difficulty forecasting sales and optimizing spend",
        "",
        "Impact:",
        "  • Lost revenue opportunities",
        "  • Inefficient marketing spend",
        "  • Poor inventory planning",
        "  • Reactive vs. proactive decisions"
    ])
    
    # Slide 3: Our Solution
    create_content_slide(prs, "Our Solution: Three-Pillar Approach", [
        "1. Data Pipeline",
        "   • Clean, process, and aggregate multi-source data",
        "   • 58 weeks of comprehensive e-commerce data",
        "",
        "2. AI Models",
        "   • Sales forecasting (89.18% accuracy)",
        "   • Adaptive learning system",
        "   • Spend optimization engine",
        "",
        "3. Interactive Dashboard",
        "   • 17 pages of real-time insights",
        "   • Production-ready deployment",
        "   • User-friendly interface"
    ])
    
    # Slide 4: Data Pipeline
    create_content_slide(prs, "Data Pipeline: 58 Weeks of Data", [
        "Data Sources:",
        "  • Traffic Overview: 1.75M visitors, 730 records",
        "  • Product Performance: 8K records",
        "  • Revenue & Sales: IDR 1.73B total",
        "  • Campaign Data: Flash Sales, Vouchers, Live, Games",
        "  • Customer Service: Chat, CSAT metrics",
        "",
        "Processing:",
        "  • Automated cleaning and validation",
        "  • Weekly aggregation for forecasting",
        "  • 25+ engineered features",
        "  • Quality checks and error handling"
    ])
    
    # Slide 5: Traffic Insights
    create_content_slide(prs, "Key Insights: Traffic Analysis", [
        "Traffic Metrics:",
        "  • Total Visitors: 1,749,117",
        "  • New Visitors: 1,410,073",
        "  • Returning Visitors: 339,044",
        "  • New vs Returning Ratio: 4.16:1",
        "",
        "Engagement:",
        "  • Average Time Spent: 1.25 minutes",
        "  • Total Products Viewed: 5.7M",
        "  • Avg Products per Visit: 3.28",
        "",
        "Insight: Strong new customer acquisition, opportunity to improve retention"
    ])
    
    # Slide 6: Sales Performance
    create_content_slide(prs, "Key Insights: Sales Performance", [
        "Revenue Metrics (58 weeks):",
        "  • Total Sales: IDR 1.73 Billion",
        "  • Average Weekly: IDR 29.87M",
        "  • Peak Week: IDR 67.87M",
        "  • Growth: +47.5% (first half vs second half)",
        "",
        "Customer Metrics:",
        "  • Total Buyers: 12,350",
        "  • Total Products Sold: 32,648",
        "  • Sales per Buyer: IDR 140,276",
        "  • Sales per Product: IDR 53,063",
        "",
        "Strong growth trajectory with healthy unit economics"
    ])
    
    # Slide 7: Campaign Performance
    create_content_slide(prs, "Key Insights: Campaign ROI", [
        "Flash Sales:",
        "  • Revenue: IDR 233M",
        "  • Orders: 1,138",
        "  • ROI: 5.2x ⭐ Best performer",
        "",
        "Vouchers:",
        "  • Cost: IDR 123M",
        "  • Orders: 1,170",
        "  • ROI: 4.5x",
        "",
        "Live Streams:",
        "  • ROI: 3.8x",
        "  • High engagement potential",
        "",
        "Recommendation: Prioritize flash sales for maximum ROI"
    ])
    
    # Slide 8: AI Model - Forecasting
    create_content_slide(prs, "AI Model #1: Sales Forecasting", [
        "XGBoost Forecasting Model",
        "",
        "Performance:",
        "  ✅ 89.18% Accuracy on test set",
        "  • MAPE: 10.82%",
        "  • R²: 0.9742",
        "  • 75% predictions within 20% error",
        "",
        "Features (25+):",
        "  • Time-based: Week, month, quarter, seasonality",
        "  • Sales patterns: Lags, rolling averages, trends",
        "  • Marketing: Ad spend, promotions, vouchers",
        "",
        "Validation: 80/20 split, 5-fold cross-validation, no overfitting"
    ])
    
    # Slide 9: Forecast Results
    create_content_slide(prs, "6-Month Sales Forecast", [
        "Predictions (Dec 2025 - Jun 2026):",
        "  • Total Forecast: IDR 0.83 Billion",
        "  • Average Weekly: IDR 31.82M",
        "  • Range: IDR 31.65M - 32.39M",
        "  • Confidence: 89.18% accuracy",
        "",
        "Business Value:",
        "  ✅ Inventory planning",
        "  ✅ Cash flow management",
        "  ✅ Staffing optimization",
        "  ✅ Marketing budget allocation",
        "",
        "40% improvement over baseline forecasting methods"
    ])
    
    # Slide 10: Feature Importance
    create_content_slide(prs, "What Drives Sales?", [
        "Top 10 Features by Importance:",
        "",
        "  1. Products Sold (37.14%) - Inventory breadth",
        "  2. Total Buyers (31.25%) - Customer acquisition",
        "  3. Product Sales (16.95%) - Revenue per product",
        "  4. Sales Trend (4.58%) - Momentum indicator",
        "  5. Total Ad Spend (2.45%) - Marketing investment",
        "  6. Has Promotion (1.89%) - Campaign activity",
        "  7. Product-Buyer Ratio (1.47%)",
        "  8. Voucher Cost (1.12%)",
        "  9. Flash Sales (0.76%)",
        "",
        "Key Insight: Product variety and customer base are primary drivers"
    ])
    
    # Slide 11: Adaptive Learning
    create_content_slide(prs, "AI Model #2: Adaptive Learning", [
        "Self-Learning Forecasting System",
        "",
        "Capability:",
        "  ✅ Automatically retrains with new weekly data",
        "  ✅ Learns from promotional campaign results",
        "  ✅ Adapts to seasonal changes",
        "  ✅ Version control and rollback",
        "",
        "How It Works:",
        "  1. New week data arrives (sales, voucher spend, etc.)",
        "  2. Model retrains on expanded dataset",
        "  3. Performance validated (85%+ accuracy threshold)",
        "  4. Model updated if quality maintained",
        "  5. Performance tracked over time",
        "",
        "Innovation: Model improves continuously without manual intervention"
    ])
    
    # Slide 12: Spend Optimizer
    create_content_slide(prs, "AI Model #3: Spend Optimizer", [
        "Budget Allocation & ROI Maximization",
        "",
        "Capability:",
        "  ✅ Test different promotional spend scenarios",
        "  ✅ Predict sales, profit, and ROI for each scenario",
        "  ✅ Compare multiple budget allocations",
        "  ✅ AI-powered recommendations",
        "",
        "Optimization Results:",
        "  • Conservative: 4.2x ROI, IDR 22M profit",
        "  • Balanced: 4.7x ROI, IDR 26M profit ⭐",
        "  • Aggressive: 4.3x ROI, IDR 28M profit",
        "",
        "Recommendation: Balanced allocation maximizes ROI efficiency"
    ])
    
    # Slide 13: Dashboard Overview
    create_content_slide(prs, "Interactive Dashboard: 17 Pages", [
        "Core Analytics:",
        "  1. Overview - KPIs and trends",
        "  2. Traffic Analysis - Visitor behavior",
        "  3. Sales Performance - Revenue metrics",
        "  4. Campaign Analytics - Promotional ROI",
        "  5. Customer Service - CSAT, response times",
        "  6. Product Recommendations - AI-powered",
        "",
        "Advanced Features:",
        "  7. Sales Forecast - 6-month predictions",
        "  8. Customer Segments - RFM analysis",
        "",
        "Innovation (NEW!):",
        "  9. Adaptive Learning - Self-improving model",
        "  10. Spend Optimizer - ROI maximization"
    ])
    
    # Slide 14: Judge Questions
    create_content_slide(prs, "Judge Questions - Answered", [
        "Q3a: Dashboard connected to Shopee API?",
        "  • No, uses historical data for demo",
        "  • Architecture ready for API integration",
        "  • Have QA credentials for production",
        "",
        "Q3b: Can it self-learn?",
        "  ✅ YES! Adaptive Learning page demonstrates this",
        "  • Automatically retrains with new weekly data",
        "  • Version control and performance tracking",
        "",
        "Q3c: Can it optimize spend?",
        "  ✅ YES! Spend Optimizer page provides this",
        "  • Test different budget scenarios",
        "  • AI-powered recommendations for ROI maximization"
    ])
    
    # Slide 15: Technical Architecture
    create_content_slide(prs, "Technical Architecture", [
        "Data Layer:",
        "  • CSV ingestion and automated cleaning",
        "  • Weekly aggregation for forecasting",
        "",
        "Model Layer:",
        "  • XGBoost forecaster (89.18% accuracy)",
        "  • Adaptive learning system",
        "  • Spend optimization engine",
        "",
        "Presentation Layer:",
        "  • Streamlit dashboard (17 pages)",
        "  • Interactive Plotly visualizations",
        "  • Real-time predictions",
        "",
        "Deployment:",
        "  • GitHub: github.com/taranggoyal70/nazava-analytics",
        "  • Cloud-ready architecture"
    ])
    
    # Slide 16: Business Impact
    create_content_slide(prs, "Business Impact & Value", [
        "Operational Efficiency:",
        "  ✅ 40% reduction in forecasting error vs. baseline",
        "  ✅ Automated weekly reporting (saves 10+ hours/week)",
        "  ✅ Real-time campaign performance tracking",
        "",
        "Revenue Optimization:",
        "  ✅ Identify optimal promotional spend (4.7x ROI)",
        "  ✅ Predict sales 6 months ahead (89% accuracy)",
        "  ✅ Reduce inventory waste through better planning",
        "",
        "Strategic Insights:",
        "  ✅ Understand key sales drivers",
        "  ✅ Optimize marketing budget allocation",
        "  ✅ Data-driven decision making"
    ])
    
    # Slide 17: Innovation Highlights
    create_content_slide(prs, "What Makes This Unique", [
        "1. Adaptive Learning:",
        "   • First Shopee analytics platform with self-learning AI",
        "   • Model improves automatically with new data",
        "   • No manual retraining required",
        "",
        "2. Spend Optimization:",
        "   • Interactive budget allocation tool",
        "   • Real-time ROI predictions",
        "   • Scenario analysis for decision support",
        "",
        "3. Comprehensive Integration:",
        "   • 12 data sources unified",
        "   • 25+ engineered features",
        "   • End-to-end analytics pipeline",
        "",
        "4. Production-Ready:",
        "   • Deployed dashboard, GitHub repository, Complete documentation"
    ])
    
    # Slide 18: Deliverables
    create_content_slide(prs, "What We Built", [
        "✅ Data Pipeline:",
        "   • 12 cleaned datasets, 58 weeks of data",
        "",
        "✅ AI Models:",
        "   • Sales forecasting (89.18% accuracy)",
        "   • Adaptive learning system",
        "   • Spend optimizer",
        "",
        "✅ Dashboard:",
        "   • 17 interactive pages",
        "   • Real-time analytics",
        "   • Production-ready",
        "",
        "✅ Documentation:",
        "   • README, SETUP guides, Jupyter notebooks",
        "",
        "✅ Deployment:",
        "   • GitHub repository, Sample data included, One-command setup"
    ])
    
    # Slide 19: Demo Instructions
    create_content_slide(prs, "Live Demo", [
        "Dashboard Access:",
        "  • URL: http://localhost:8501",
        "  • Login: admin / admin123",
        "  • GitHub: github.com/taranggoyal70/nazava-analytics",
        "",
        "Demo Flow:",
        "  1. Overview - Show KPIs and trends",
        "  2. Sales Forecast - 6-month prediction",
        "  3. Adaptive Learning - Add new data, retrain model",
        "  4. Spend Optimizer - Test budget scenarios",
        "  5. Campaign Analytics - ROI analysis",
        "",
        "Try It Yourself:",
        "  git clone https://github.com/taranggoyal70/nazava-analytics",
        "  cd nazava-analytics/dashboard",
        "  streamlit run app.py"
    ])
    
    # Slide 20: Conclusion
    create_content_slide(prs, "Conclusion", [
        "Problem:",
        "  • Shopee sellers lack unified analytics and forecasting tools",
        "",
        "Solution:",
        "  • Nazava Analytics Platform with AI-powered insights",
        "",
        "Results:",
        "  • 89.18% forecasting accuracy",
        "  • 4.7x optimal ROI identified",
        "  • Self-learning model capability",
        "  • 17-page interactive dashboard",
        "",
        "Impact:",
        "  • Better inventory planning",
        "  • Optimized marketing spend",
        "  • Data-driven decisions",
        "  • Increased profitability"
    ])
    
    # Slide 21: Thank You
    create_title_slide(
        prs,
        "Thank You!",
        "Questions?\n\nGitHub: github.com/taranggoyal70/nazava-analytics\nDashboard: http://localhost:8501\nLogin: admin / admin123"
    )
    
    # Save presentation
    filename = "Nazava_Analytics_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"\nSlides created: {len(prs.slides)}")
    print("\nOpen the file to view your presentation!")

if __name__ == "__main__":
    main()
